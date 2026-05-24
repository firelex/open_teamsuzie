from __future__ import annotations

from pathlib import Path

import pytest

from pptx_template_agent.injection import (
    compute_capacity,
    fill_deck,
    inspect_template,
    lint_pptx,
    will_overflow,
)
from pptx_template_agent.models import DeckSpec, SlideUpdate

TEMPLATE = Path(__file__).parent.parent / "templates" / "ic-template-v3.pptx"
pytestmark = pytest.mark.skipif(not TEMPLATE.exists(), reason="IC template not present")


def test_capacity_embedded_in_manifest():
    m = inspect_template(TEMPLATE)
    # Find a known sub-header text frame
    found = False
    for slide in m["slides"]:
        for s in slide["shapes"]:
            if s["kind"] == "text" and s["capacity"]:
                cap = s["capacity"]
                assert cap["max_chars_per_line"] > 0
                assert cap["max_lines"] > 0
                assert cap["total_chars"] == cap["max_chars_per_line"] * cap["max_lines"]
                found = True
    assert found, "expected at least one text frame with computed capacity"


def test_will_overflow_long_header():
    cap = {
        "font_pt": 28.0, "width_in": 11.3, "height_in": 1.0,
        "max_chars_per_line": 58, "max_lines": 2, "total_chars": 116,
    }
    # 117 chars across 1 paragraph wraps to 3 lines → overflow
    assert will_overflow("x" * 117, cap) is True
    # 100 chars wraps to 2 lines → fits
    assert will_overflow("x" * 100, cap) is False
    # multi-paragraph: 3 paragraphs of 30 chars each = 3 lines → overflow
    assert will_overflow(["x" * 30, "x" * 30, "x" * 30], cap) is True


def test_filler_flags_overflow(tmp_path: Path):
    """An 180-char header should overflow the ~116-char sub-header capacity."""
    long_header = "Exec summary | " + ("a very long marketing line " * 6).strip()
    spec = DeckSpec(
        deal_name="X",
        template=str(TEMPLATE),
        slides=[SlideUpdate(
            slide_index=2,
            fields={"Text Placeholder 2": long_header},
        )],
    )
    _, report = fill_deck(spec, tmp_path / "out.pptx")
    assert report.overflows, "expected at least one overflow warning"
    assert "slide2:Text Placeholder 2" in report.overflows[0]


def test_filler_no_overflow_for_short_header(tmp_path: Path):
    spec = DeckSpec(
        deal_name="X",
        template=str(TEMPLATE),
        slides=[SlideUpdate(
            slide_index=2,
            fields={"Text Placeholder 2": "Exec summary | Short"},
        )],
    )
    _, report = fill_deck(spec, tmp_path / "out.pptx")
    assert report.overflows == []


def test_linter_finds_unfilled_stock_patterns(tmp_path: Path):
    """A deck where most slides are skipped should surface unfilled stock patterns."""
    spec = DeckSpec(
        deal_name="X",
        template=str(TEMPLATE),
        slides=[SlideUpdate(slide_index=0, fields={"Text Placeholder 1": "Cover"})],
    )
    out, report = fill_deck(spec, tmp_path / "out.pptx")
    # The template has many [xxx] / [Role] / xxx tokens in unfilled slides
    assert len(report.unfilled) > 5
    # Spot-check a known stock pattern (e.g. "xxx" in a table cell)
    samples = [w["text"] for w in report.unfilled]
    assert any("xxx" in s.lower() or "[" in s for s in samples)


def test_effective_width_never_exceeds_raw_width():
    """The effective width must always be ≤ the shape's raw width — this
    catches accidental sign flips or overflow in the intrusion logic."""
    from pptx import Presentation as _P
    prs = _P(str(TEMPLATE))
    from pptx_template_agent.injection.capacity import _effective_width
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame or not shape.width:
                continue
            eff = _effective_width(shape, slide)
            assert eff <= shape.width, (
                f"{shape.name} on slide: eff={eff} > raw={shape.width}"
            )
            assert eff > 0


def test_lint_cell_overflows_flags_oversize_text(tmp_path: Path):
    """A table cell whose content exceeds (cols × rows × font) should be
    surfaced by lint_cell_overflows."""
    from pptx_template_agent.injection import lint_cell_overflows
    # Write a deliberately-long string into a known-narrow cell on slide 5
    # (VCP table). The VCP table column 0 ("VC lever") is narrow.
    spec = DeckSpec(
        deal_name="Test",
        template=str(TEMPLATE),
        slides=[SlideUpdate(
            slide_index=5,
            tables={"Tabelle 3": [
                ["VC lever", "T", "U", "Today", "D", "A"],
                ["A" * 200, "x", "x", "x", "x", "x"],  # very long content in narrow col
            ]},
        )],
    )
    path, report = fill_deck(spec, tmp_path / "out.pptx")
    # The 200-char value should trigger a cell-overflow warning
    assert report.cell_overflows, "expected at least one cell overflow"
    assert any(
        w["row"] == 1 and w["col"] == 0 for w in report.cell_overflows
    )


def test_lint_table_collisions_uses_estimated_bottom(tmp_path: Path):
    """The collision detector should account for per-row natural heights,
    not just the table's declared bounding box. On the Forecast slide
    (Table 12 with 16 rows at 0.19" declared but 11pt text needing more),
    the rendered bottom pushes past Rectangle 5 ("Takeaways") — bounding-box
    analysis misses it, estimated bottom catches it."""
    from pptx_template_agent.injection import lint_table_collisions
    import sys as _sys
    _sys.path.insert(0, str(Path(__file__).parent.parent / "examples"))
    from aquapure import build  # type: ignore[import-not-found]

    spec = build()
    # Detection only — disable the auto-fix so the saved file still has the
    # underlying collision when linted.
    out, _ = fill_deck(spec, tmp_path / "out.pptx", auto_fix_collisions=False)
    collisions = lint_table_collisions(out)
    forecast_collision = next(
        (c for c in collisions
         if c["slide_index"] == 14 and c["table_name"] == "Table 12"
         and c["overlaps_with"] == "Rectangle 5"),
        None,
    )
    assert forecast_collision is not None, (
        "expected to detect Forecast Table 12 vs Rectangle 5 collision"
    )


def test_auto_resolve_collisions_eliminates_forecast_overlap(tmp_path: Path):
    """With auto-fix enabled (default), the Forecast Table 12 vs Rectangle 5
    collision must be resolved at fill time and absent from the lint pass."""
    from pptx_template_agent.injection import lint_table_collisions
    import sys as _sys
    _sys.path.insert(0, str(Path(__file__).parent.parent / "examples"))
    from aquapure import build  # type: ignore[import-not-found]

    spec = build()
    out, report = fill_deck(spec, tmp_path / "out.pptx")  # auto_fix on
    assert report.collision_fixes > 0
    collisions = lint_table_collisions(out)
    # No collision should remain between Forecast Table 12 and Rectangle 5
    forecast_collision = next(
        (c for c in collisions
         if c["slide_index"] == 14 and c["table_name"] == "Table 12"
         and c["overlaps_with"] == "Rectangle 5"),
        None,
    )
    assert forecast_collision is None


def test_lint_table_collisions_runs_cleanly(tmp_path: Path):
    """The collision detector should run on every saved file without raising
    and return a list (possibly empty for clean layouts)."""
    from pptx_template_agent.injection import lint_table_collisions
    spec = DeckSpec(
        deal_name="Test",
        template=str(TEMPLATE),
        slides=[SlideUpdate(slide_index=0, fields={"Text Placeholder 1": "X"})],
    )
    path, _ = fill_deck(spec, tmp_path / "out.pptx")
    warnings = lint_table_collisions(path)
    assert isinstance(warnings, list)
    # Each entry must carry the expected keys
    for w in warnings[:3]:
        assert {"slide_index", "table_name", "overlaps_with",
                "table_bottom_in", "other_top_in", "overlap_in"}.issubset(w.keys())


def test_lint_pptx_standalone(tmp_path: Path):
    """lint_pptx can be called on any .pptx, not just our outputs."""
    warnings = lint_pptx(TEMPLATE)
    assert len(warnings) > 10  # the template is full of placeholders
    # Every entry has the expected shape
    for w in warnings[:3]:
        assert {"slide_index", "shape_name", "text", "pattern"}.issubset(w.keys())

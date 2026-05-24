from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from pptx_template_agent.injection import fill_deck, inspect_template
from pptx_template_agent.models import DeckSpec, SlideUpdate

TEMPLATE = Path(__file__).parent.parent / "templates" / "ic-template-v3.pptx"
pytestmark = pytest.mark.skipif(not TEMPLATE.exists(), reason="IC template not present")


def test_inspect_template_returns_manifest():
    m = inspect_template(TEMPLATE)
    assert m["slide_count"] >= 1
    assert m["slide_width_in"] > 0
    assert m["slides"][0]["layout"]
    assert any(s["kind"] == "table" for slide in m["slides"] for s in slide["shapes"])


def test_fill_deck_writes_pptx_and_applies_text(tmp_path: Path):
    spec = DeckSpec(
        deal_name="Test",
        template=str(TEMPLATE),
        slides=[SlideUpdate(slide_index=0, fields={
            "Text Placeholder 1": "Project X",
            "Text Placeholder 2": "Today",
        })],
    )
    out = tmp_path / "out.pptx"
    path, report = fill_deck(spec, out)
    assert path.exists()
    assert report.fields_applied == 2
    assert report.fields_missing == []

    prs = Presentation(str(path))
    texts = [s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame]
    assert "Project X" in texts
    assert "Today" in texts


def test_fill_deck_reports_missing_shapes(tmp_path: Path):
    spec = DeckSpec(
        deal_name="Test",
        template=str(TEMPLATE),
        slides=[SlideUpdate(slide_index=0, fields={"NoSuchShape": "x"})],
    )
    _, report = fill_deck(spec, tmp_path / "out.pptx")
    assert report.fields_applied == 0
    assert "slide0:NoSuchShape" in report.fields_missing


def test_fill_deck_raises_on_bad_slide_index(tmp_path: Path):
    spec = DeckSpec(
        deal_name="Test",
        template=str(TEMPLATE),
        slides=[SlideUpdate(slide_index=999, fields={})],
    )
    with pytest.raises(IndexError):
        fill_deck(spec, tmp_path / "out.pptx")


def test_list_content_creates_multiple_paragraphs(tmp_path: Path):
    spec = DeckSpec(
        deal_name="Test",
        template=str(TEMPLATE),
        slides=[SlideUpdate(slide_index=0, fields={
            "Text Placeholder 2": ["line one", "line two", "line three"],
        })],
    )
    path, _ = fill_deck(spec, tmp_path / "out.pptx")
    prs = Presentation(str(path))
    shape = next(s for s in prs.slides[0].shapes if s.name == "Text Placeholder 2")
    assert len(shape.text_frame.paragraphs) == 3
    assert shape.text_frame.paragraphs[1].text == "line two"

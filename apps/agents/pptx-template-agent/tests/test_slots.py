from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from pptx_template_agent.injection import fill_deck, inspect_template
from pptx_template_agent.models import DeckSpec, SlideUpdate

TEMPLATE = Path(__file__).parent.parent / "templates" / "ic-template-v3.pptx"
pytestmark = pytest.mark.skipif(not TEMPLATE.exists(), reason="IC template not present")


def test_slot_detection_on_exec_summary():
    """The exec summary slide has 4 ovals (1/2/3/4) over Content Placeholder 2.
    Only that placeholder should report slots."""
    m = inspect_template(TEMPLATE)
    slide = m["slides"][2]
    slot_shapes = [s for s in slide["shapes"] if s["kind"] == "text" and s.get("slots")]
    assert len(slot_shapes) == 1
    shape = slot_shapes[0]
    assert shape["name"] == "Content Placeholder 2"
    info = shape["slots"]
    assert info["count"] == 4
    assert info["labels"] == ["1", "2", "3", "4"]


def test_no_false_slot_detection_on_short_shapes():
    """Small shapes (sub-headers, ovals themselves) must not get slot info."""
    m = inspect_template(TEMPLATE)
    slide = m["slides"][2]
    for s in slide["shapes"]:
        if s["name"].startswith("Oval"):
            assert s.get("slots") is None, f"Badge {s['name']} should not have slots"
        if s["name"] == "Text Placeholder 4":
            # the sub-header on this slide is 0.24" tall — too short for slots
            assert s.get("slots") is None


def test_slot_fill_writes_section_starts_with_spcbef(tmp_path: Path):
    """A list-of-lists content for a slot-enabled frame should emit
    content-only paragraphs and set explicit `<a:spcBef>` on the first
    paragraph of each section after the first — pushing it down to
    align with its slot anchor without blank-paragraph padding."""
    from pptx.oxml.ns import qn

    content = [
        ["A1", "A2"],
        ["B1"],
        ["C1", "C2"],
        ["D1"],
    ]
    spec = DeckSpec(
        deal_name="Slots",
        template=str(TEMPLATE),
        slides=[SlideUpdate(
            slide_index=2,
            fields={"Content Placeholder 2": content},
        )],
    )
    path, report = fill_deck(spec, tmp_path / "out.pptx")
    assert report.fields_applied >= 1

    prs = Presentation(str(path))
    shape = next(
        s for s in prs.slides[2].shapes if s.name == "Content Placeholder 2"
    )
    paras = list(shape.text_frame.paragraphs)
    texts = [p.text for p in paras]
    # All section content should appear. No blank padding paragraphs.
    assert texts == ["A1", "A2", "B1", "C1", "C2", "D1"]

    def spc_bef_val(p) -> int | None:
        pPr = p._p.find(qn("a:pPr"))
        if pPr is None:
            return None
        spc = pPr.find(qn("a:spcBef"))
        if spc is None:
            return None
        pts = spc.find(qn("a:spcPts"))
        if pts is None or pts.get("val") is None:
            return None
        return int(pts.get("val"))

    # First section (index 0, 1) — no override expected.
    # Section starts (index 2, 3, 5) — explicit spcBef set.
    assert spc_bef_val(paras[2]) is not None and spc_bef_val(paras[2]) > 0
    assert spc_bef_val(paras[3]) is not None and spc_bef_val(paras[3]) > 0
    assert spc_bef_val(paras[5]) is not None and spc_bef_val(paras[5]) > 0

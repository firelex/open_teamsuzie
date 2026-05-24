from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from pptx_template_agent.injection import (
    analyse_table,
    fill_deck,
    inspect_template,
    normalize_cell,
    normalize_row,
)
from pptx_template_agent.models import DeckSpec, NormalizationOp, SlideUpdate

TEMPLATE = Path(__file__).parent.parent / "templates" / "ic-template-v3.pptx"
pytestmark = pytest.mark.skipif(not TEMPLATE.exists(), reason="IC template not present")


def test_analyse_table_classifies_rows():
    prs = Presentation(str(TEMPLATE))
    tbl = next(s.table for s in prs.slides[10].shapes if s.has_table)
    analysis = analyse_table(tbl)
    assert analysis["body_modal_pt"] == 10.0
    kinds = [r["kind"] for r in analysis["rows"]]
    # Row 0 must be header; rows 1, 7, 11, 12 are dividers/empty in this template
    assert kinds[0] == "header"
    assert kinds[1] == "section_divider"
    assert kinds[7] == "section_divider"
    assert kinds[12] == "section_divider"
    # And several should be body
    assert kinds.count("body") >= 8


def test_manifest_embeds_row_schema_and_anomalies():
    m = inspect_template(TEMPLATE)
    # Find a table shape — slide 10 historicals
    table_shape = next(
        s for slide in m["slides"]
        if slide["slide_index"] == 10
        for s in slide["shapes"] if s["kind"] == "table"
    )
    assert table_shape["body_modal_pt"] == 10.0
    assert len(table_shape["row_schema"]) > 0
    # Clean template should have no anomalies on this table
    assert table_shape["anomalies"] == []


def test_normalize_cell(tmp_path: Path):
    """Snap a cell's font size to a different pt and verify it stuck."""
    spec = DeckSpec(
        deal_name="N",
        template=str(TEMPLATE),
        slides=[SlideUpdate(
            slide_index=10,
            normalize={"Table 7": [NormalizationOp(row=2, col=1, target_pt=14.0)]},
            tables={"Table 7": [["", "", "", "", "", ""], ["", "", "", "", "", ""], ["Revenue", "9999"]]},
        )],
    )
    path, report = fill_deck(spec, tmp_path / "out.pptx")
    assert report.normalisations_applied >= 1

    prs = Presentation(str(path))
    tbl = next(s.table for s in prs.slides[10].shapes if s.has_table)
    cell = tbl.cell(2, 1)
    sizes = [
        run.font.size.pt
        for para in cell.text_frame.paragraphs
        for run in para.runs
        if run.font.size is not None
    ]
    assert 14.0 in sizes


def test_normalize_row_applies_to_every_column(tmp_path: Path):
    """Whole-row normalize hits every column."""
    spec = DeckSpec(
        deal_name="N",
        template=str(TEMPLATE),
        slides=[SlideUpdate(
            slide_index=10,
            normalize={"Table 7": [NormalizationOp(row=2, target_pt=12.0)]},
        )],
    )
    path, report = fill_deck(spec, tmp_path / "out.pptx")
    assert report.normalisations_applied == 6  # 6 columns

    prs = Presentation(str(path))
    tbl = next(s.table for s in prs.slides[10].shapes if s.has_table)
    for c in range(6):
        runs = list(tbl.cell(2, c).text_frame.paragraphs[0].runs)
        for run in runs:
            if run.font.size is not None:
                assert run.font.size.pt == 12.0


def test_filler_assigns_explicit_pt_to_cells_lacking_one(tmp_path: Path):
    """Cells whose template counterpart had no run-level size must get the
    table's body modal pt after fill — this is the defensive fix that
    prevents 'oversized D&A row' style bugs."""
    spec = DeckSpec(
        deal_name="X",
        template=str(TEMPLATE),
        slides=[SlideUpdate(
            slide_index=10,
            tables={"Table 7": [
                ["h1", "h2", "h3", "h4", "h5", "h6"],
                # row 1 in the template is a section divider; cells 1-5 had no run-level size
                ["P&L", "x", "x", "x", "x", "x"],
            ]},
        )],
    )
    path, _ = fill_deck(spec, tmp_path / "out.pptx")
    prs = Presentation(str(path))
    tbl = next(s.table for s in prs.slides[10].shapes if s.has_table)
    # The previously-empty cells in row 1 should now carry an explicit pt
    for c in range(1, 6):
        sizes = [
            run.font.size.pt
            for para in tbl.cell(1, c).text_frame.paragraphs
            for run in para.runs
            if run.font.size is not None
        ]
        assert sizes, f"cell (1,{c}) lacks explicit pt after fill"
        assert sizes[0] == 10.0  # the table's body modal

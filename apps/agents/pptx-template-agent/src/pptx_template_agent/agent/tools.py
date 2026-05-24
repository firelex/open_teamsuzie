"""Anthropic-style tool definitions and dispatch for the agent loop.

Each tool wraps an injection-layer primitive. State lives in an `AgentState`
instance that the loop creates per job — this keeps the tool functions pure
in terms of inputs and easy to test."""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation

from ..injection import inspect_template
from ..injection.capacity import compute_cell_capacity, will_overflow
from ..injection.filler import fill_deck
from ..models import DeckSpec, NormalizationOp, SlideUpdate
from ..services.templates import get_template_path

TOOL_DEFINITIONS: list[dict[str, Any]] = [
    {
        "name": "inspect_template",
        "description": (
            "Read the template's structure: slide count, layout names, shape names "
            "per slide, current placeholder text, and table dimensions. Always call "
            "this first to discover what you can fill."
        ),
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": [],
        },
    },
    {
        "name": "update_slide",
        "description": (
            "Queue an update to a slide. Field values can be:\n"
            "  - a single string (one paragraph)\n"
            "  - a list of strings (multiple paragraphs / bullets)\n"
            "  - a list-of-lists when the target frame has `slots` in the "
            "manifest — one inner list per slot anchor, in order; provide "
            "exactly as many inner lists as `slots.count`.\n"
            "Tables are 2D matrices applied top-left aligned. RAG values: "
            "green|amber|red|grey applied to a shape name or 'TableName!r,c'. "
            "Use `normalize` to snap a table cell or row to an explicit pt "
            "BEFORE content is written.\n\n"
            "Pre-flight capacity check: the dispatcher checks every field "
            "and table cell against the target shape's capacity. If anything "
            "overflows, the result includes a `warnings` list and the update "
            "is NOT queued — shorten the content and call again."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "slide_index": {"type": "integer"},
                "fields": {
                    "type": "object",
                    "additionalProperties": {
                        "oneOf": [
                            {"type": "string"},
                            {"type": "array", "items": {"type": "string"}},
                            {
                                "type": "array",
                                "description": (
                                    "Slot-aware form: one inner list per slot "
                                    "anchor (use when manifest reports slots)"
                                ),
                                "items": {
                                    "type": "array",
                                    "items": {"type": "string"},
                                },
                            },
                        ]
                    },
                },
                "tables": {
                    "type": "object",
                    "additionalProperties": {
                        "type": "array",
                        "items": {"type": "array", "items": {"type": "string"}},
                    },
                },
                "rag": {
                    "type": "object",
                    "additionalProperties": {
                        "type": "string",
                        "enum": ["green", "amber", "red", "grey"],
                    },
                },
                "normalize": {
                    "type": "object",
                    "description": "table_name -> list of normalisation ops",
                    "additionalProperties": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "row": {"type": "integer"},
                                "col": {"type": "integer"},
                                "target_pt": {"type": "number"},
                            },
                            "required": ["row", "target_pt"],
                        },
                    },
                },
            },
            "required": ["slide_index"],
        },
    },
    {
        "name": "save_deck",
        "description": (
            "Write the accumulated updates to a .pptx file and stop. Call this "
            "exactly once at the end of the run. The response includes any "
            "overflows / cell_overflows / unfilled warnings the filler "
            "produced."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "deal_name": {"type": "string"},
            },
            "required": ["deal_name"],
        },
    },
]


# ---- AgentState --------------------------------------------------------------


@dataclass
class AgentState:
    template_path: Path
    output_path: Path
    updates: list[SlideUpdate] = field(default_factory=list)
    saved: bool = False
    save_report: dict[str, Any] | None = None

    # Cached on first use — used for pre-flight capacity validation
    _manifest: Any = None
    _presentation: Any = None

    def manifest(self) -> Any:
        if self._manifest is None:
            self._manifest = inspect_template(self.template_path)
        return self._manifest

    def presentation(self) -> Any:
        """Open the raw .pptx (read-only) so we can probe live shape /
        table objects when validating cell capacity. Cached for the run."""
        if self._presentation is None:
            self._presentation = Presentation(str(self.template_path))
        return self._presentation


# ---- Pre-flight validation ---------------------------------------------------


def _shape_in_manifest(manifest, slide_index: int, shape_name: str):
    if slide_index >= len(manifest["slides"]):
        return None
    for s in manifest["slides"][slide_index]["shapes"]:
        if s["name"] == shape_name:
            return s
    return None


def _shape_capacity_in_manifest(manifest, slide_index: int, shape_name: str):
    shape = _shape_in_manifest(manifest, slide_index, shape_name)
    return shape.get("capacity") if shape else None


def _table_dims_in_manifest(manifest, slide_index: int, table_name: str):
    if slide_index >= len(manifest["slides"]):
        return None
    for s in manifest["slides"][slide_index]["shapes"]:
        if s["name"] == table_name and s["kind"] == "table":
            return s
    return None


def _find_table(prs, slide_index: int, table_name: str):
    if slide_index >= len(prs.slides):
        return None
    for shape in prs.slides[slide_index].shapes:
        if shape.name == table_name and shape.has_table:
            return shape.table
    return None


def _is_nested_list(value: Any) -> bool:
    return isinstance(value, list) and len(value) > 0 and isinstance(value[0], list)


def _flatten_for_overflow_check(content: Any) -> Any:
    """For overflow checks, treat slot-mode list-of-lists as a flat paragraph
    list (all sections concatenated). The slot-aware filler interleaves
    blank-line padding, but content length is the sum of all sections."""
    if _is_nested_list(content):
        flat: list[str] = []
        for sub in content:
            if isinstance(sub, list):
                flat.extend(sub)
            else:
                flat.append(str(sub))
        return flat
    return content


def _validate_update_capacity(state: AgentState, slide_index: int,
                               fields: dict, tables: dict) -> list[dict]:
    """Return a list of warning dicts for any field / cell that would
    overflow its target or violate slot rules. Empty list = good to go."""
    warnings: list[dict] = []
    manifest = state.manifest()

    # Field-level checks
    for name, content in (fields or {}).items():
        shape = _shape_in_manifest(manifest, slide_index, name)
        if shape is None:
            # Unknown shape — let the filler's missing-field report catch it
            continue
        cap = shape.get("capacity")
        slot_info = shape.get("slots")

        # Slot enforcement: if the frame has N slot anchors, the content
        # MUST be a list-of-lists with exactly N non-empty inner lists.
        if slot_info and slot_info.get("count", 0) >= 2:
            expected = slot_info["count"]
            if not _is_nested_list(content):
                warnings.append({
                    "kind": "slot_shape_required",
                    "slide_index": slide_index,
                    "shape": name,
                    "expected_inner_lists": expected,
                    "guidance": (
                        f"This frame has {expected} slot anchors. Pass `fields[\"{name}\"]` "
                        f"as a list with exactly {expected} inner lists — one per slot — "
                        "each containing at least one non-empty paragraph."
                    ),
                })
            else:
                if len(content) != expected:
                    warnings.append({
                        "kind": "slot_count_mismatch",
                        "slide_index": slide_index,
                        "shape": name,
                        "expected_inner_lists": expected,
                        "got": len(content),
                    })
                else:
                    empty_slots = [
                        i for i, sub in enumerate(content)
                        if not any(str(p).strip() for p in (sub if isinstance(sub, list) else [sub]))
                    ]
                    if empty_slots:
                        warnings.append({
                            "kind": "empty_slot",
                            "slide_index": slide_index,
                            "shape": name,
                            "empty_slot_indices": empty_slots,
                            "guidance": (
                                "Every slot anchor needs at least one non-empty paragraph. "
                                f"Fill slots {empty_slots} with content or shorten the outline "
                                "and pass fewer slots overall."
                            ),
                        })

        if cap is None:
            continue
        if will_overflow(_flatten_for_overflow_check(content), cap):
            if isinstance(content, str):
                sample = content
            elif _is_nested_list(content):
                sample = " | ".join(
                    " ".join(str(p) for p in sub) for sub in content if sub
                )
            else:
                sample = " | ".join(str(p) for p in content)
            warnings.append({
                "kind": "field_overflow",
                "slide_index": slide_index,
                "shape": name,
                "max_chars_per_line": cap["max_chars_per_line"],
                "max_lines": cap["max_lines"],
                "text_preview": sample[:80],
            })

    # Cell-level checks (column width × row height)
    prs = state.presentation() if tables else None
    for table_name, matrix in (tables or {}).items():
        if prs is None:
            continue
        table = _find_table(prs, slide_index, table_name)
        if table is None:
            continue  # filler reports unknown tables
        for r, row in enumerate(matrix):
            if r >= len(table.rows):
                break
            for c, value in enumerate(row):
                if c >= len(table.columns):
                    break
                if not value or not str(value).strip():
                    continue
                cap = compute_cell_capacity(table, r, c)
                if cap is None:
                    continue
                if will_overflow(str(value), cap):
                    warnings.append({
                        "kind": "cell_overflow",
                        "slide_index": slide_index,
                        "table": table_name,
                        "row": r, "col": c,
                        "max_chars_per_line": cap["max_chars_per_line"],
                        "max_lines": cap["max_lines"],
                        "text_preview": str(value)[:60],
                    })

    return warnings


# ---- Dispatch ----------------------------------------------------------------


def dispatch(state: AgentState, tool_name: str, tool_input: dict[str, Any]) -> Any:
    if tool_name == "inspect_template":
        return state.manifest()

    if tool_name == "update_slide":
        slide_index = tool_input["slide_index"]
        fields = tool_input.get("fields", {})
        tables = tool_input.get("tables", {})

        # Pre-flight capacity validation — refuse to queue updates that would
        # produce visible overflow. The model must shorten and re-call.
        warnings = _validate_update_capacity(state, slide_index, fields, tables)
        if warnings:
            return {
                "error": "capacity overflow — update not queued",
                "warnings": warnings,
                "guidance": (
                    "Shorten the content of each flagged field/cell to fit "
                    "within max_chars_per_line × max_lines, then call "
                    "update_slide again with the same slide_index."
                ),
            }

        raw_normalize = tool_input.get("normalize") or {}
        normalize = {
            t: [NormalizationOp(**op) for op in ops]
            for t, ops in raw_normalize.items()
        }
        update = SlideUpdate(
            slide_index=slide_index,
            fields=fields,
            tables=tables,
            rag=tool_input.get("rag", {}),
            normalize=normalize,
        )
        state.updates.append(update)
        return {"ok": True, "queued_updates": len(state.updates)}

    if tool_name == "save_deck":
        # Compute which slides the agent touched, so the post-save scan
        # can flag *unfilled placeholders on slides the agent worked on*
        # (vs slides the agent intentionally skipped, where leftover
        # template defaults are expected).
        touched_slides = {u.slide_index for u in state.updates}

        spec = DeckSpec(
            deal_name=tool_input["deal_name"],
            template=str(state.template_path),
            slides=state.updates,
        )
        path, report = fill_deck(spec, state.output_path)

        unfilled_on_touched = [
            dict(w) for w in report.unfilled
            if w["slide_index"] in touched_slides
        ]

        state.saved = True
        state.save_report = {
            "path": str(path),
            "fields_applied": report.fields_applied,
            "tables_applied": report.tables_applied,
            "rag_applied": report.rag_applied,
            "missing_fields": report.fields_missing,
            "missing_tables": report.tables_missing,
            "missing_rag": report.rag_missing,
            "overflows": report.overflows,
            "unfilled": [dict(w) for w in report.unfilled],
            "unfilled_on_touched_slides": unfilled_on_touched,
            "cell_overflows": [dict(w) for w in report.cell_overflows],
            "table_collisions": [dict(w) for w in report.table_collisions],
            "cells_shrunk": report.cells_shrunk,
            "collision_fixes": report.collision_fixes,
        }
        return state.save_report

    raise ValueError(f"Unknown tool: {tool_name}")


def state_from_template_id(template_id: str, output_path: Path) -> AgentState:
    return AgentState(
        template_path=get_template_path(template_id),
        output_path=output_path,
    )

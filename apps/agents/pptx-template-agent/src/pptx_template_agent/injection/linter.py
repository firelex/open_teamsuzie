"""Post-fill linter: scan the saved .pptx for template stock patterns that
the spec didn't override, e.g. '[xxx]', '[xxxx]', '[Placeholder ...]'.

The agent uses these warnings to decide whether to backfill, clear the
shape, or accept the default text."""

from __future__ import annotations

import re
from pathlib import Path
from typing import TypedDict

from pptx import Presentation

from .capacity import compute_cell_capacity, will_overflow

# Stock patterns we treat as "unfilled" — case-insensitive, anchored on the
# bracketed/placeholder shapes the IC template uses. Tweak conservatively;
# false positives are noisy but never destructive (the linter only warns).
_PATTERNS = [
    re.compile(r"^\[?xxx+\]?$", re.IGNORECASE),                # xxx, [xxx], [xxxx], XXX
    re.compile(r"\[\s*placeholder", re.IGNORECASE),            # [Placeholder — ...]
    re.compile(r"^\[(role|name|thesis|x|xxx+)[^\]]*\]$", re.IGNORECASE),
    re.compile(r"^\[[^\]]{1,40}\]$"),                          # any short bracketed token left in-place
]


class UnfilledWarning(TypedDict):
    slide_index: int
    shape_name: str
    text: str
    pattern: str


class CellOverflowWarning(TypedDict):
    slide_index: int
    table_name: str
    row: int
    col: int
    text: str
    max_chars_per_line: int
    max_lines: int


class TableCollisionWarning(TypedDict):
    slide_index: int
    table_name: str
    overlaps_with: str
    table_bottom_in: float
    other_top_in: float
    overlap_in: float


def _is_stock(text: str) -> str | None:
    stripped = text.strip()
    if not stripped:
        return None
    for pat in _PATTERNS:
        if pat.search(stripped):
            return pat.pattern
    return None


def lint_pptx(path: str | Path) -> list[UnfilledWarning]:
    """Walk every text frame + table cell in the .pptx and report any text
    that still matches an "unfilled" stock pattern."""
    prs = Presentation(str(path))
    warnings: list[UnfilledWarning] = []

    def visit_text_frame(slide_index: int, shape_name: str, tf) -> None:
        for para in tf.paragraphs:
            txt = para.text
            pat = _is_stock(txt)
            if pat:
                warnings.append(UnfilledWarning(
                    slide_index=slide_index, shape_name=shape_name,
                    text=txt[:80], pattern=pat,
                ))

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                visit_text_frame(i, shape.name, shape.text_frame)
            if shape.has_table:
                for r, row in enumerate(shape.table.rows):
                    for c, cell in enumerate(row.cells):
                        visit_text_frame(i, f"{shape.name}!{r},{c}", cell.text_frame)

    return warnings


def _line_height_pt(font_pt: float) -> float:
    return font_pt * 1.2


def _row_height_pt(table, row: int) -> float:
    h = table.rows[row].height or 0
    return h / 12700  # EMU → pt


# PowerPoint applies ~5pt of additional vertical padding per cell (≈ 0.07")
# beyond the run's line height. Empirically calibrated against an 11pt font
# in a 0.19" declared row rendering at ~0.26" — gives natural row height
# 11*1.2 + 5 = 18.2pt ≈ 0.25". Use when estimating *rendered* row height.
_CELL_VPAD_PT = 5.0


def _estimated_row_height_pt(table, row: int) -> float:
    """Estimate the row's *rendered* height: max of the declared height and
    the largest cell's natural text height (font × 1.2 + cell padding)."""
    declared_pt = _row_height_pt(table, row)
    max_font_pt = 0.0
    for c in range(len(table.columns)):
        cap = compute_cell_capacity(table, row, c)
        if cap and cap["font_pt"] > max_font_pt:
            max_font_pt = cap["font_pt"]
    if max_font_pt == 0:
        return declared_pt
    natural_pt = _line_height_pt(max_font_pt) + _CELL_VPAD_PT
    return max(declared_pt, natural_pt)


def _estimated_table_bottom_emu(shape) -> int:
    """Effective bottom of a table accounting for per-row natural heights —
    catches tables whose declared height is shorter than their content
    actually needs, which is the root cause of "table overlaps the shape
    below it" in PowerPoint's own rendering."""
    table = shape.table
    y = shape.top
    for r in range(len(table.rows)):
        y += int(_estimated_row_height_pt(table, r) * 12700)
    return y


def lint_cell_overflows(path: str | Path) -> list[CellOverflowWarning]:
    """Scan every populated table cell and flag those that visually overflow:
    either the text exceeds (width × height × font), OR the font's line height
    is larger than the row height (so even a single character won't fit
    cleanly in the row's vertical space — common when a template author
    bumped a row's font but the row height wasn't grown to match).
    """
    prs = Presentation(str(path))
    warnings: list[CellOverflowWarning] = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            for r in range(len(table.rows)):
                row_h_pt = _row_height_pt(table, r)
                for c in range(len(table.columns)):
                    cell = table.cell(r, c)
                    text = cell.text.strip()
                    if not text:
                        continue
                    cap = compute_cell_capacity(table, r, c)
                    if cap is None:
                        continue
                    overflows = will_overflow(text, cap)
                    font_too_tall = (
                        row_h_pt > 0
                        and _line_height_pt(cap["font_pt"]) > row_h_pt * 1.1
                    )
                    if overflows or font_too_tall:
                        warnings.append(CellOverflowWarning(
                            slide_index=i,
                            table_name=shape.name,
                            row=r, col=c,
                            text=text[:80],
                            max_chars_per_line=cap["max_chars_per_line"],
                            max_lines=cap["max_lines"],
                        ))
    return warnings


def lint_table_collisions(path: str | Path) -> list[TableCollisionWarning]:
    """For each table on each slide, compute the table's *rendered* bottom
    (per-row natural heights, not just the declared bounding box) and flag
    any shape whose top sits within that effective range.

    PowerPoint renders rows at max(declared_height, font_line_height + ~2pt
    padding) — when the template's declared row height is tight, the rendered
    table extends past its declared box, colliding with shapes positioned
    immediately below it. The estimated bottom catches this; the bounding-box
    bottom misses it."""
    prs = Presentation(str(path))
    warnings: list[TableCollisionWarning] = []
    EMU_PER_INCH = 914400
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_table or not shape.width or not shape.height:
                continue
            t_top = shape.top
            t_bottom = _estimated_table_bottom_emu(shape)
            t_left = shape.left
            t_right = shape.left + shape.width
            for other in slide.shapes:
                if other is shape or not other.width or not other.height:
                    continue
                o_top = other.top
                if o_top <= t_top:
                    continue                       # not below the table top
                if o_top >= t_bottom:
                    continue                       # below the table bottom — no collision
                # require horizontal overlap to be a real visual collision
                o_left = other.left
                o_right = other.left + other.width
                if o_right <= t_left or o_left >= t_right:
                    continue
                # skip empty text frames — they're typically inactive
                # sub-header placeholders positioned mid-slide that are
                # never visually rendered when blank
                if other.has_text_frame and not other.text_frame.text.strip():
                    continue
                warnings.append(TableCollisionWarning(
                    slide_index=i,
                    table_name=shape.name,
                    overlaps_with=other.name,
                    table_bottom_in=t_bottom / EMU_PER_INCH,
                    other_top_in=o_top / EMU_PER_INCH,
                    overlap_in=(t_bottom - o_top) / EMU_PER_INCH,
                ))
    return warnings

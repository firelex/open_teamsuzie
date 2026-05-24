"""Template introspection: walk an arbitrary .pptx and emit a manifest of
slide indices, layout names, and the shape names that can be written to.

This is the primitive that lets the agent (and API clients) work against any
.pptx template without hard-coded knowledge of its structure.
"""

from __future__ import annotations

from pathlib import Path
from typing import Literal, TypedDict

from pptx import Presentation

from .anomalies import Anomaly, RowStyle, analyse_table
from .capacity import TextCapacity, compute_capacity
from .slots import SlotInfo, detect_slots


class ShapeManifest(TypedDict):
    name: str
    kind: Literal["text", "table", "chart", "picture", "group", "other"]
    sample_text: str  # current text frame content, truncated — useful for LLMs
    rows: int  # tables only; 0 otherwise
    cols: int  # tables only; 0 otherwise
    # tables only: row-by-row classification + body modal pt + anomalies
    body_modal_pt: float | None
    row_schema: list[RowStyle]
    anomalies: list[Anomaly]
    # text frames only: how much content fits before overflow
    capacity: TextCapacity | None
    # text frames only: numbered-badge anchors that subdivide the frame
    slots: SlotInfo | None


class SlideManifest(TypedDict):
    slide_index: int
    layout: str
    shapes: list[ShapeManifest]


class TemplateManifest(TypedDict):
    slide_width_in: float
    slide_height_in: float
    slide_count: int
    slides: list[SlideManifest]


def _shape_kind(shape) -> str:
    if shape.has_table:
        return "table"
    if shape.has_chart:
        return "chart"
    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
        return "picture"
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        return "group"
    if shape.has_text_frame:
        return "text"
    return "other"


def inspect_template(template_path: str | Path, *, sample_chars: int = 120) -> TemplateManifest:
    """Open the .pptx at `template_path` and return its manifest.

    Reading the file does not modify it. The returned manifest is JSON-safe.
    """
    path = Path(template_path)
    if not path.exists():
        raise FileNotFoundError(f"Template not found: {path}")

    prs = Presentation(str(path))
    slides: list[SlideManifest] = []
    for idx, slide in enumerate(prs.slides):
        shapes: list[ShapeManifest] = []
        for shape in slide.shapes:
            kind = _shape_kind(shape)
            sample = ""
            if shape.has_text_frame:
                sample = shape.text_frame.text.replace("\n", " | ")[:sample_chars]
            rows = cols = 0
            body_modal_pt: float | None = None
            row_schema: list[RowStyle] = []
            anomalies: list[Anomaly] = []
            if shape.has_table:
                rows = len(shape.table.rows)
                cols = len(shape.table.columns)
                analysis = analyse_table(shape.table)
                body_modal_pt = analysis["body_modal_pt"]
                row_schema = analysis["rows"]
                anomalies = analysis["anomalies"]
            capacity = compute_capacity(shape, slide) if kind == "text" else None
            slot_info: SlotInfo | None = None
            if kind == "text":
                slot_info = detect_slots(
                    slide, shape, font_pt=capacity["font_pt"] if capacity else None,
                )
            shapes.append(
                ShapeManifest(
                    name=shape.name,
                    kind=kind,  # type: ignore[arg-type]
                    sample_text=sample,
                    rows=rows,
                    cols=cols,
                    body_modal_pt=body_modal_pt,
                    row_schema=row_schema,
                    anomalies=anomalies,
                    capacity=capacity,
                    slots=slot_info,
                )
            )
        slides.append(
            SlideManifest(
                slide_index=idx,
                layout=slide.slide_layout.name,
                shapes=shapes,
            )
        )

    return TemplateManifest(
        slide_width_in=prs.slide_width / 914400,
        slide_height_in=prs.slide_height / 914400,
        slide_count=len(prs.slides),
        slides=slides,
    )

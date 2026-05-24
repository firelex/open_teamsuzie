"""High-level deck filler: takes a DeckSpec, writes a .pptx.

The filler is template-agnostic — `spec.template` may point at any .pptx.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation

from ..models import DeckSpec, SlideUpdate
from .capacity import compute_capacity, will_overflow
from .linter import (
    CellOverflowWarning,
    TableCollisionWarning,
    UnfilledWarning,
    lint_cell_overflows,
    lint_pptx,
    lint_table_collisions,
)
from .normalize import auto_shrink_oversized_fonts, normalize_cell, normalize_row
from .template_fixup import auto_resolve_collisions
from .placeholders import set_text_by_name
from .tables import apply_rag, set_table

log = logging.getLogger(__name__)


@dataclass
class FillReport:
    fields_applied: int = 0
    fields_missing: list[str] = field(default_factory=list)
    tables_applied: int = 0
    tables_missing: list[str] = field(default_factory=list)
    rag_applied: int = 0
    rag_missing: list[str] = field(default_factory=list)
    normalisations_applied: int = 0
    overflows: list[str] = field(default_factory=list)
    unfilled: list[UnfilledWarning] = field(default_factory=list)
    cell_overflows: list[CellOverflowWarning] = field(default_factory=list)
    table_collisions: list[TableCollisionWarning] = field(default_factory=list)
    cells_shrunk: int = 0
    collision_fixes: int = 0

    @property
    def is_clean(self) -> bool:
        return not (
            self.fields_missing or self.tables_missing or self.rag_missing
            or self.overflows or self.unfilled
            or self.cell_overflows or self.table_collisions
        )


def _apply_slide(prs: Presentation, update: SlideUpdate, report: FillReport) -> None:
    if update.slide_index >= len(prs.slides):
        raise IndexError(
            f"slide_index {update.slide_index} out of range "
            f"(template has {len(prs.slides)} slides)"
        )
    slide = prs.slides[update.slide_index]

    # Apply normalisations FIRST so subsequent table writes inherit the
    # corrected font sizes (set_table picks up the new modal automatically).
    for table_name, ops in update.normalize.items():
        for op in ops:
            if op.col is None:
                applied = normalize_row(slide, table_name, op.row, op.target_pt)
            else:
                applied = 1 if normalize_cell(
                    slide, table_name, op.row, op.col, op.target_pt
                ) else 0
            report.normalisations_applied += applied

    for name, content in update.fields.items():
        # Capacity pre-check — flag overflows but still write the content
        # (truncation would lose information silently; the warning is enough).
        for shape in slide.shapes:
            if shape.name == name and shape.has_text_frame:
                cap = compute_capacity(shape, slide)
                if cap and will_overflow(content, cap):
                    sample = (content if isinstance(content, str) else " | ".join(content))[:60]
                    report.overflows.append(
                        f"slide{update.slide_index}:{name} "
                        f"(max~{cap['max_chars_per_line']}×{cap['max_lines']}): {sample!r}"
                    )
                break

        n = set_text_by_name(slide, name, content)
        if n == 0:
            report.fields_missing.append(f"slide{update.slide_index}:{name}")
        else:
            report.fields_applied += n

    for table_name, matrix in update.tables.items():
        if set_table(slide, table_name, matrix):
            report.tables_applied += 1
        else:
            report.tables_missing.append(f"slide{update.slide_index}:{table_name}")

    for target, rag in update.rag.items():
        n = apply_rag(slide, target, rag)
        if n == 0:
            report.rag_missing.append(f"slide{update.slide_index}:{target}")
        else:
            report.rag_applied += n


def fill_deck(
    spec: DeckSpec,
    output_path: str | Path,
    *,
    auto_shrink_cells: bool = True,
    auto_fix_collisions: bool = True,
) -> tuple[Path, FillReport]:
    """Open the template referenced by `spec.template`, apply the spec, write
    to `output_path`. Returns (path, report).

    When `auto_shrink_cells` is True (default), the filler walks every table
    cell after the spec is applied and snaps any cell whose explicit font
    size would make its line height exceed the row height. Catches template
    authoring artifacts (e.g. someone bumped a row's font but never resized
    the row) without requiring the spec to know about them."""
    template_path = Path(spec.template)
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    prs = Presentation(str(template_path))
    report = FillReport()

    # Run template fixup *before* applying the spec so downstream-shape
    # positions reflect the corrected layout in any post-fill checks.
    if auto_fix_collisions:
        report.collision_fixes = len(auto_resolve_collisions(prs))

    for update in spec.slides:
        _apply_slide(prs, update, report)

    if auto_shrink_cells:
        report.cells_shrunk = auto_shrink_oversized_fonts(prs)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    # Post-fill linters scan the saved file for issues the agent should know
    # about — stock patterns left unfilled, cell content that overflows its
    # column×row capacity, and tables that visually collide with adjacent shapes.
    report.unfilled = lint_pptx(out)
    report.cell_overflows = lint_cell_overflows(out)
    report.table_collisions = lint_table_collisions(out)

    log.info(
        "Wrote %s (fields=%d tables=%d rag=%d missing=%d/%d/%d "
        "overflows=%d unfilled=%d cell_overflows=%d collisions=%d "
        "shrunk=%d collision_fixes=%d)",
        out, report.fields_applied, report.tables_applied, report.rag_applied,
        len(report.fields_missing), len(report.tables_missing), len(report.rag_missing),
        len(report.overflows), len(report.unfilled),
        len(report.cell_overflows), len(report.table_collisions),
        report.cells_shrunk, report.collision_fixes,
    )
    return out, report

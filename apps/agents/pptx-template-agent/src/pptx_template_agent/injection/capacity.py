"""Text-frame capacity: how many characters / lines fit before overflow.

Uses a heuristic: average sans-serif char width ≈ 0.5 × font_size_pt. That's
within ~10% of PowerPoint's actual rendering for the brand fonts we see in
practice (Inter, Arial, Calibri, Open Sans). Good enough for "will this
overflow?" — refine to PIL font metrics if you need pixel accuracy later.

EMU constants — PowerPoint's internal unit:
  1 inch     = 914400 EMU
  1 point    = 12700 EMU      (72 pt per inch)
"""

from __future__ import annotations

from typing import TypedDict

from pptx.oxml.ns import qn

EMU_PER_PT = 12700
EMU_PER_INCH = 914400
AVG_CHAR_WIDTH_RATIO = 0.5      # sans-serif heuristic
# Fallback line-spacing ratio used only when the text frame's OOXML
# inheritance chain doesn't carry an explicit <a:lnSpc>. The resolver
# below walks run → paragraph → text-frame defaults → layout placeholder
# → master and returns the template-declared spacing whenever it exists.
LINE_HEIGHT_RATIO_FALLBACK = 1.2
DEFAULT_FONT_PT = 10.0          # fallback when we can't read the font size


class TextCapacity(TypedDict):
    font_pt: float
    width_in: float
    height_in: float
    max_chars_per_line: int
    max_lines: int
    total_chars: int            # max_chars_per_line × max_lines (rough cap)


def _sz_to_pt(sz_attr: str | None) -> float | None:
    """OOXML font sizes are stored as 100×pt, e.g. sz='2800' == 28pt."""
    if not sz_attr:
        return None
    try:
        return float(sz_attr) / 100.0
    except (TypeError, ValueError):
        return None


def _xml_defRPr_size(element) -> float | None:
    """Walk an element's descendants for the first <a:defRPr sz="..."> hit."""
    if element is None:
        return None
    for def_rpr in element.iter(qn("a:defRPr")):
        pt = _sz_to_pt(def_rpr.get("sz"))
        if pt is not None:
            return pt
    return None


def _ln_spc_value(ln_spc_el, font_pt: float) -> float | None:
    """Resolve an <a:lnSpc> element to a line-height in pt.

    OOXML supports two forms:
      <a:spcPct val="120000"/>  → 120% of the font size (multiplier × 100000)
      <a:spcPts val="1800"/>    → 18pt exactly (value × 100)
    """
    if ln_spc_el is None:
        return None
    spc_pct = ln_spc_el.find(qn("a:spcPct"))
    if spc_pct is not None:
        val = spc_pct.get("val")
        try:
            return font_pt * (int(val) / 100000.0)
        except (TypeError, ValueError):
            return None
    spc_pts = ln_spc_el.find(qn("a:spcPts"))
    if spc_pts is not None:
        val = spc_pts.get("val")
        try:
            return int(val) / 100.0
        except (TypeError, ValueError):
            return None
    return None


def _xml_spacing_pt(element, tag: str) -> float | None:
    """Find the first <a:spcBef> or <a:spcAft> child and return its
    spcPts value in pt (or None when only spcPct is used / nothing
    declared).

    `tag` should be 'a:spcBef' or 'a:spcAft'."""
    if element is None:
        return None
    for spc in element.iter(qn(tag)):
        pts = spc.find(qn("a:spcPts"))
        if pts is not None:
            val = pts.get("val")
            try:
                return int(val) / 100.0
            except (TypeError, ValueError):
                continue
    return None


def _resolve_paragraph_spacing(shape) -> tuple[float, float]:
    """Walk the same inheritance chain we use for line spacing and return
    (spc_bef_pt, spc_aft_pt). Each paragraph in this frame is rendered
    with `spc_bef` above and `spc_aft` below. Falls back to (0, 0) when
    neither is declared anywhere."""
    if not shape.has_text_frame:
        return (0.0, 0.0)
    tf = shape.text_frame

    def _hunt(tag: str) -> float | None:
        # Per-paragraph override on the first paragraph (highest priority —
        # `set_shape_text` clones this pPr when emitting new paragraphs, so
        # this is what new paragraphs will inherit unless we override).
        for p in tf.paragraphs:
            pPr = p._p.find(qn("a:pPr"))
            if pPr is not None:
                v = _xml_spacing_pt(pPr, tag)
                if v is not None:
                    return v
            break
        body_el = tf._txBody
        for tag_name in ("a:lstStyle", "a:bodyPr"):
            sub = body_el.find(qn(tag_name))
            v = _xml_spacing_pt(sub, tag)
            if v is not None:
                return v
        # Matching layout placeholder
        if shape.is_placeholder:
            try:
                idx = shape.placeholder_format.idx
                layout = shape.part.slide_layout
            except AttributeError:
                layout = None
            if layout is not None:
                for sp in layout.element.iter(qn("p:sp")):
                    ph = sp.find(".//" + qn("p:ph"))
                    if ph is None or str(ph.get("idx") or 0) != str(idx):
                        continue
                    txbody = sp.find(qn("p:txBody"))
                    v = _xml_spacing_pt(txbody, tag)
                    if v is not None:
                        return v
                    break
        # Master text styles
        try:
            master = shape.part.slide_layout.slide_master
        except AttributeError:
            master = None
        if master is not None:
            tx_styles = master.element.find(qn("p:txStyles"))
            root = tx_styles if tx_styles is not None else master.element
            v = _xml_spacing_pt(root, tag)
            if v is not None:
                return v
        return None

    return (_hunt("a:spcBef") or 0.0, _hunt("a:spcAft") or 0.0)


def _resolve_line_height_pt(shape, font_pt: float) -> float:
    """Walk the OOXML inheritance chain for the first <a:lnSpc> that applies
    to this text frame's paragraphs. Falls back to font_pt × ratio when no
    explicit spacing is declared."""
    if not shape.has_text_frame:
        return font_pt * LINE_HEIGHT_RATIO_FALLBACK
    tf = shape.text_frame

    # 1. Run / paragraph-level explicit spacing (rare but valid)
    body_el = tf._txBody
    for tag in ("a:lstStyle", "a:bodyPr"):
        sub = body_el.find(qn(tag))
        if sub is None:
            continue
        ln = next(sub.iter(qn("a:lnSpc")), None)
        h = _ln_spc_value(ln, font_pt)
        if h is not None:
            return h

    # 2. Matching layout placeholder
    if shape.is_placeholder:
        try:
            idx = shape.placeholder_format.idx
            layout = shape.part.slide_layout
        except AttributeError:
            layout = None
        if layout is not None:
            for sp in layout.element.iter(qn("p:sp")):
                ph = sp.find(".//" + qn("p:ph"))
                if ph is None or str(ph.get("idx") or 0) != str(idx):
                    continue
                txbody = sp.find(qn("p:txBody"))
                if txbody is None:
                    continue
                ln = next(txbody.iter(qn("a:lnSpc")), None)
                h = _ln_spc_value(ln, font_pt)
                if h is not None:
                    return h
                break

    # 3. Master text styles
    try:
        master = shape.part.slide_layout.slide_master
    except AttributeError:
        master = None
    if master is not None:
        tx_styles = master.element.find(qn("p:txStyles"))
        root = tx_styles if tx_styles is not None else master.element
        ln = next(root.iter(qn("a:lnSpc")), None)
        h = _ln_spc_value(ln, font_pt)
        if h is not None:
            return h

    return font_pt * LINE_HEIGHT_RATIO_FALLBACK


def _layout_placeholder_size(shape) -> float | None:
    """Find the matching placeholder on the slide layout and read its
    lstStyle/lvl1pPr/defRPr sz. This is the step the master fallback misses."""
    if not shape.is_placeholder:
        return None
    try:
        idx = shape.placeholder_format.idx
        layout = shape.part.slide_layout
    except AttributeError:
        return None

    for sp in layout.element.iter(qn("p:sp")):
        ph = sp.find(".//" + qn("p:ph"))
        if ph is None:
            continue
        if str(ph.get("idx") or 0) != str(idx):
            continue
        txbody = sp.find(qn("p:txBody"))
        if txbody is None:
            return None
        return _xml_defRPr_size(txbody)
    return None


def _master_style_size(shape) -> float | None:
    """Resolve the master text style for a placeholder shape (titleStyle /
    bodyStyle / otherStyle) and return the lvl1 defRPr font size in pt."""
    if not shape.is_placeholder:
        return None
    try:
        ph_type = shape.placeholder_format.type  # MSO_PLACEHOLDER enum
    except Exception:
        return None
    # Walk up to the slide master (SlidePart.slide_layout.slide_master)
    try:
        master = shape.part.slide_layout.slide_master
    except AttributeError:
        return None

    # bodyStyle / titleStyle / otherStyle are inside p:txStyles
    tx_styles = master.element.find(qn("p:txStyles"))
    master_el = tx_styles if tx_styles is not None else master.element
    # Title-family vs body-family vs other
    title_kinds = {13, 15, 14}   # TITLE, CENTER_TITLE, SUBTITLE — best-effort
    body_kinds = {2, 7, 18}      # BODY, OBJECT, BODY-ish
    if ph_type and getattr(ph_type, "value", ph_type) in title_kinds:
        style_tag = "p:titleStyle"
    elif ph_type and getattr(ph_type, "value", ph_type) in body_kinds:
        style_tag = "p:bodyStyle"
    else:
        style_tag = "p:bodyStyle"  # bodyStyle is the closest universal fallback

    style_el = master_el.find(qn(style_tag))
    return _xml_defRPr_size(style_el)


def _first_run_font_pt(shape) -> float | None:
    """Best-effort font-size resolution for the shape's text frame, walking
    run → paragraph → text_frame defaults → master text style."""
    if not shape.has_text_frame:
        return None
    tf = shape.text_frame

    # 1. Run-level explicit size
    for para in tf.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                return run.font.size.pt

    # 2. Paragraph-level explicit size
    for para in tf.paragraphs:
        if para.font.size is not None:
            return para.font.size.pt

    # 3. Text frame's lstStyle / bodyPr defRPr
    body_el = tf._txBody
    for tag in ("a:lstStyle", "a:bodyPr"):
        sub = body_el.find(qn(tag))
        pt = _xml_defRPr_size(sub)
        if pt is not None:
            return pt

    # 4. Matching placeholder on the slide layout (lstStyle defRPr)
    pt = _layout_placeholder_size(shape)
    if pt is not None:
        return pt

    # 5. Master's title/body/otherStyle (lvl1 defRPr)
    pt = _master_style_size(shape)
    if pt is not None:
        return pt

    return None


def _effective_width(shape, slide=None) -> int:
    """If another shape's bounding box overlaps this shape's vertical band
    and starts to the right of this shape's left edge, the usable width is
    reduced to the gap. Catches headers whose right edge is constrained by
    an adjacent box (e.g. a "Comments" rail on the right of a slide)."""
    width_emu = shape.width or 0
    if width_emu <= 0 or slide is None:
        return width_emu
    s_left = shape.left
    s_right = shape.left + shape.width
    s_top = shape.top
    s_bottom = shape.top + shape.height
    smallest_intrusion_left = s_right
    for other in slide.shapes:
        if other is shape:
            continue
        if not other.width or not other.height:
            continue
        o_left = other.left
        o_top = other.top
        o_bottom = other.top + other.height
        # vertical overlap must be substantial — at least 50% of the smaller
        # shape's height. Catches "comments box beside a header" but skips
        # icons sitting just above/below a card whose edge barely crosses.
        overlap_top = max(s_top, o_top)
        overlap_bottom = min(s_bottom, o_bottom)
        v_overlap = overlap_bottom - overlap_top
        if v_overlap <= 0:
            continue
        min_h = min(s_bottom - s_top, o_bottom - o_top)
        if v_overlap < min_h * 0.5:
            continue
        # Must start meaningfully to the right (≥ 0.2") — co-located shapes
        # are typically containers/backgrounds, not horizontal intrusions.
        if o_left - s_left < int(0.2 * EMU_PER_INCH):
            continue
        # Intruder must not extend past us on the right — if its right edge is
        # >= our right, it's likely a container/wrapper, not a horizontal cut.
        if other.left + other.width >= s_right:
            continue
        smallest_intrusion_left = min(smallest_intrusion_left, o_left)
    return max(1, smallest_intrusion_left - s_left)


def compute_capacity(shape, slide=None) -> TextCapacity | None:
    """Compute capacity for a shape's text frame. Returns None if the shape
    has no text frame or no usable width/height.

    Pass `slide` to enable spatial-overlap analysis — the effective width
    accounts for adjacent shapes that constrain the text frame's usable area.
    """
    if not shape.has_text_frame:
        return None
    height_emu = shape.height or 0
    if not shape.width or shape.width <= 0 or height_emu <= 0:
        return None

    width_emu = _effective_width(shape, slide)

    font_pt = _first_run_font_pt(shape) or DEFAULT_FONT_PT
    line_height_pt = _resolve_line_height_pt(shape, font_pt)
    char_width_pt = font_pt * AVG_CHAR_WIDTH_RATIO
    char_width_emu = char_width_pt * EMU_PER_PT
    line_height_emu = line_height_pt * EMU_PER_PT

    max_chars_per_line = max(1, int(width_emu / char_width_emu))
    max_lines = max(1, int(height_emu / line_height_emu))

    return TextCapacity(
        font_pt=font_pt,
        width_in=width_emu / 914400,
        height_in=height_emu / 914400,
        max_chars_per_line=max_chars_per_line,
        max_lines=max_lines,
        total_chars=max_chars_per_line * max_lines,
    )


def compute_cell_capacity(table, row: int, col: int) -> TextCapacity | None:
    """Capacity for a single table cell, using the column width and row height.

    Cells inherit the same font-resolution chain as standalone text frames —
    we walk run → paragraph → text-frame defaults. Master fallback is skipped
    because cells don't have a placeholder type."""
    try:
        cell = table.cell(row, col)
    except (IndexError, ValueError):
        return None

    col_width_emu = table.columns[col].width or 0
    row_height_emu = table.rows[row].height or 0
    if col_width_emu <= 0 or row_height_emu <= 0:
        return None

    # Use the same run/paragraph/text-frame resolution as standalone shapes.
    # cell.text_frame has the same interface as shape.text_frame.
    tf = cell.text_frame
    font_pt: float | None = None
    for para in tf.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                font_pt = run.font.size.pt
                break
        if font_pt:
            break
    if font_pt is None:
        for para in tf.paragraphs:
            if para.font.size is not None:
                font_pt = para.font.size.pt
                break
    if font_pt is None:
        for tag in ("a:lstStyle", "a:bodyPr"):
            sub = tf._txBody.find(qn(tag))
            font_pt = _xml_defRPr_size(sub)
            if font_pt is not None:
                break
    if font_pt is None:
        font_pt = DEFAULT_FONT_PT

    # Resolve line spacing from the cell's own text frame. Cells don't
    # inherit from placeholders/masters the way standalone shapes do, so
    # we only look at the cell's local <a:lnSpc> and fall back to the
    # ratio if none is declared.
    ln = None
    for tag in ("a:lstStyle", "a:bodyPr"):
        sub = tf._txBody.find(qn(tag))
        if sub is None:
            continue
        ln = next(sub.iter(qn("a:lnSpc")), None)
        if ln is not None:
            break
    cell_line_pt = _ln_spc_value(ln, font_pt)
    line_height_pt = cell_line_pt if cell_line_pt is not None else font_pt * LINE_HEIGHT_RATIO_FALLBACK

    char_width_emu = font_pt * AVG_CHAR_WIDTH_RATIO * EMU_PER_PT
    line_height_emu = line_height_pt * EMU_PER_PT

    max_chars_per_line = max(1, int(col_width_emu / char_width_emu))
    max_lines = max(1, int(row_height_emu / line_height_emu))

    return TextCapacity(
        font_pt=font_pt,
        width_in=col_width_emu / EMU_PER_INCH,
        height_in=row_height_emu / EMU_PER_INCH,
        max_chars_per_line=max_chars_per_line,
        max_lines=max_lines,
        total_chars=max_chars_per_line * max_lines,
    )


def will_overflow(content: str | list[str], capacity: TextCapacity) -> bool:
    """Quick check: does `content` overflow the capacity?

    Counts hard newlines as line breaks and assumes each paragraph wraps
    independently. Approximate but conservative — flags more often than
    PowerPoint would actually overflow."""
    if isinstance(content, str):
        paragraphs = content.split("\n")
    else:
        paragraphs = list(content)

    used_lines = 0
    for para in paragraphs:
        if not para:
            used_lines += 1
            continue
        # ceil-divide para length by max_chars_per_line for the wrap count
        n = (len(para) + capacity["max_chars_per_line"] - 1) // capacity["max_chars_per_line"]
        used_lines += max(1, n)
    return used_lines > capacity["max_lines"]

"""Explicit normalisation of table cell font sizes.

The set_table filler already sets a defensive fallback pt on every cell.
Use these helpers when you want to override a specific cell or row to a
chosen pt — for example, to "snap" a template anomaly back to the row mode
that the agent decided was the intent."""

from __future__ import annotations

from pptx.slide import Slide
from pptx.util import Pt

from .tables import find_table_by_name


def normalize_cell(slide: Slide, table_name: str, row: int, col: int, target_pt: float) -> bool:
    table = find_table_by_name(slide, table_name)
    if table is None or row >= len(table.rows) or col >= len(table.columns):
        return False
    cell = table.cell(row, col)
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(target_pt)
    return True


def normalize_row(slide: Slide, table_name: str, row: int, target_pt: float) -> int:
    table = find_table_by_name(slide, table_name)
    if table is None or row >= len(table.rows):
        return 0
    count = 0
    for c in range(len(table.columns)):
        if normalize_cell(slide, table_name, row, c, target_pt):
            count += 1
    return count


def auto_shrink_oversized_fonts(presentation, *, line_height_ratio: float = 1.2) -> int:
    """Walk every table on every slide. For each cell whose explicit font
    size implies a line height larger than the cell's row height, snap the
    font down to (row_height_pt / line_height_ratio).

    Returns the number of cells adjusted. Template-agnostic: any oversized
    style left in a template (or set by a template author who bumped fonts
    without resizing rows) gets normalised to fit the row.
    """
    adjusted = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            for r in range(len(table.rows)):
                row_h_emu = table.rows[r].height or 0
                if row_h_emu <= 0:
                    continue
                row_h_pt = row_h_emu / 12700
                target_pt = row_h_pt / line_height_ratio
                for c in range(len(table.columns)):
                    cell = table.cell(r, c)
                    changed = False
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size is None:
                                continue
                            current_pt = run.font.size.pt
                            line_h_pt = current_pt * line_height_ratio
                            if line_h_pt > row_h_pt * 1.1:
                                run.font.size = Pt(target_pt)
                                changed = True
                    if changed:
                        adjusted += 1
    return adjusted

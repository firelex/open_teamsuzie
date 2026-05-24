"""Template fixup: detect table-shape collisions and resolve them by
shifting downstream shapes down. Designed to run once at ingestion so
every subsequent use of the template sees a clean layout — the LLM agent
authoring against the manifest no longer has to think about it."""

from __future__ import annotations

from typing import TypedDict

from pptx.util import Pt

from .linter import _CELL_VPAD_PT, _estimated_table_bottom_emu

EMU_PER_INCH = 914400
EMU_PER_PT = 12700


class CollisionFix(TypedDict):
    slide_index: int
    table_name: str
    moved_shape: str
    shift_in: float


def _is_anchored(shape) -> bool:
    """Shapes we deliberately never move: slide-number placeholders are
    typically anchored to a fixed slide footer location. Moving them
    breaks the brand layout."""
    return "Slide Number" in shape.name


def _shrink_table_to_fit(table_shape, available_height_emu: int) -> None:
    """Shrink every cell's explicit font in the table so the natural rendered
    height of the table fits within `available_height_emu`.

    Computes the per-row pt budget from the available height divided by the
    number of rows, then snaps every run's font.size to (budget - cell vpad)
    / 1.2. Only shrinks — never grows."""
    table = table_shape.table
    n_rows = len(table.rows)
    if n_rows == 0 or available_height_emu <= 0:
        return
    per_row_pt = (available_height_emu / EMU_PER_PT) / n_rows
    target_pt = max(6.0, (per_row_pt - _CELL_VPAD_PT) / 1.2)
    for r in range(n_rows):
        for c in range(len(table.columns)):
            cell = table.cell(r, c)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is None:
                        continue
                    if run.font.size.pt > target_pt:
                        run.font.size = Pt(target_pt)


def auto_resolve_collisions(
    presentation,
    *,
    buffer_in: float = 0.05,
) -> list[CollisionFix]:
    """For each table that overflows its declared bottom into another
    shape's territory, shift that shape (and every shape positioned at or
    below it) down by the overlap amount + buffer.

    Caps the shift so no shape's bottom is pushed past the slide bottom —
    a smaller-than-ideal shift is preferred over off-slide content.

    Operates in-place on the Presentation. Returns the list of fixes
    applied. Idempotent: a second call on the same presentation returns []."""
    fixes: list[CollisionFix] = []
    buffer_emu = int(buffer_in * EMU_PER_INCH)
    slide_height_emu = presentation.slide_height

    for i, slide in enumerate(presentation.slides):
        for table_shape in [s for s in slide.shapes if s.has_table]:
            t_top = table_shape.top
            t_left = table_shape.left
            t_right = table_shape.left + table_shape.width
            t_bottom = _estimated_table_bottom_emu(table_shape)

            # Find the topmost shape that collides with the table's rendered area
            colliders = []
            for other in slide.shapes:
                if other is table_shape:
                    continue
                if not other.width or not other.height:
                    continue
                if other.top <= t_top:
                    continue                       # above or aligned with table top
                if other.top >= t_bottom:
                    continue                       # already below the rendered table
                # require real horizontal overlap
                o_left = other.left
                o_right = other.left + other.width
                if o_right <= t_left or o_left >= t_right:
                    continue
                # skip empty placeholders and anchored shapes
                if other.has_text_frame and not other.text_frame.text.strip():
                    continue
                if _is_anchored(other):
                    continue
                colliders.append(other)

            if not colliders:
                continue

            colliders.sort(key=lambda s: s.top)
            top_collider = colliders[0]
            new_top = t_bottom + buffer_emu
            desired_shift = new_top - top_collider.top
            if desired_shift <= 0:
                continue

            # Find the lowest-bottom shape we'd be moving — cap the shift so
            # its bottom stays within the slide.
            cutoff = top_collider.top
            movable = [
                s for s in slide.shapes
                if s is not table_shape
                and s.top is not None and s.top >= cutoff
                and s.height
                and not _is_anchored(s)
            ]
            if movable:
                max_bottom = max(s.top + s.height for s in movable)
                headroom = slide_height_emu - max_bottom
                shift = min(desired_shift, max(0, headroom))
            else:
                shift = desired_shift

            for other in movable:
                if shift > 0:
                    other.top = other.top + shift
                    fixes.append(CollisionFix(
                        slide_index=i,
                        table_name=table_shape.name,
                        moved_shape=other.name,
                        shift_in=shift / EMU_PER_INCH,
                    ))

            # If the cap left residual overlap, shrink the table's fonts so
            # its rendered height fits in the (now larger) available space.
            residual_overlap = desired_shift - shift
            if residual_overlap > 0:
                new_collider_top = top_collider.top  # after the shift above
                available = new_collider_top - table_shape.top - buffer_emu
                _shrink_table_to_fit(table_shape, available)
                fixes.append(CollisionFix(
                    slide_index=i,
                    table_name=table_shape.name,
                    moved_shape="<font shrink>",
                    shift_in=0.0,
                ))

    return fixes

"""Set text on shapes by shape name, preserving the template's typography.

python-pptx replaces a `text_frame.text = "..."` by wiping the runs and creating
a fresh default run — which loses the template's font, size, color, weight.
The functions here instead reuse the first run's formatting as the style
template for every new paragraph, so the slide master's design survives.
"""

from __future__ import annotations

from copy import deepcopy

from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.text.text import _Paragraph, _Run

from .slots import detect_slots


def _clone_run_style(source_run: _Run, target_run: _Run) -> None:
    """Copy run-level rPr (font, size, color, bold) from source to target."""
    src_rpr = source_run._r.find(qn("a:rPr"))
    if src_rpr is None:
        return
    # Drop any existing rPr on target first
    existing = target_run._r.find(qn("a:rPr"))
    if existing is not None:
        target_run._r.remove(existing)
    new_rpr = deepcopy(src_rpr)
    target_run._r.insert(0, new_rpr)


def _set_paragraph_spc_bef(p_element, spc_bef_units: int | None) -> None:
    """Set or clear the `<a:spcBef>` value on a paragraph's pPr.

    `spc_bef_units` is OOXML's pt × 100. None removes any existing
    spcBef override (lets the paragraph inherit). Idempotent — replaces
    any prior override on this paragraph."""
    pPr = p_element.find(qn("a:pPr"))
    if pPr is None:
        if spc_bef_units is None:
            return
        pPr = p_element.makeelement(qn("a:pPr"), {})
        p_element.insert(0, pPr)
    # Drop any existing spcBef so we don't end up with two
    for existing in pPr.findall(qn("a:spcBef")):
        pPr.remove(existing)
    if spc_bef_units is None:
        return
    spc_bef = pPr.makeelement(qn("a:spcBef"), {})
    spc_pts = spc_bef.makeelement(qn("a:spcPts"), {"val": str(max(0, int(spc_bef_units)))})
    spc_bef.append(spc_pts)
    # spcBef must appear in the correct schema order — prepend keeps it
    # near the top of pPr where the schema expects it.
    pPr.insert(0, spc_bef)


def _set_paragraph_text(paragraph: _Paragraph, text: str, style_run: _Run | None) -> None:
    """Replace all runs in `paragraph` with a single run holding `text`,
    cloning `style_run`'s formatting if provided."""
    # Remove every <a:r> child
    for r in paragraph._p.findall(qn("a:r")):
        paragraph._p.remove(r)
    # Also clean up line breaks (<a:br>) — we'll redo them via multi-line.
    for br in paragraph._p.findall(qn("a:br")):
        paragraph._p.remove(br)

    new_run = paragraph.add_run()
    new_run.text = text
    if style_run is not None:
        _clone_run_style(style_run, new_run)


def _style_run_from(text_frame) -> _Run | None:
    """Find a representative run in the text frame to clone style from."""
    for para in text_frame.paragraphs:
        for run in para.runs:
            return run
    return None


def _slot_aware_paragraphs(shape, content: list[list[str]]) -> list[tuple[str, int | None]]:
    """Convert a list-of-lists slot spec into an ordered list of
    `(text, spc_bef_emu)` tuples — content paragraphs only, no blank
    padding. The spc_bef_emu is set on the first paragraph of each
    section after the first to push it down to its slot anchor exactly;
    other paragraphs leave spc_bef inherited (None).

    PowerPoint must honour explicit `<a:spcBef>` regardless of font /
    line-height assumptions, so this is the deterministic path.
    Returns an empty list if no slots were detected (caller falls back
    to flat join)."""
    from .capacity import (
        compute_capacity,
        _resolve_line_height_pt,
        _resolve_paragraph_spacing,
        _first_run_font_pt,
    )

    slide = shape.part.slide
    slots = detect_slots(slide, shape)
    if not slots or slots["count"] < 2:
        return []

    font_pt = _first_run_font_pt(shape) or 12.0
    line_h_pt = _resolve_line_height_pt(shape, font_pt)
    spc_bef_pt, spc_aft_pt = _resolve_paragraph_spacing(shape)
    cap = compute_capacity(shape, slide)
    max_chars = cap["max_chars_per_line"] if cap else 60

    def visual_lines(text: str) -> int:
        if not text:
            return 1
        return max(1, (len(text) + max_chars - 1) // max_chars)

    # For each section after the first, compute an spcBef override that
    # makes the section's first paragraph TOP land at its slot anchor.
    # We emit EXPLICIT spc_bef on every paragraph (override or default in
    # pt × 100 OOXML units) so cloned pPrs can't propagate the override
    # to subsequent paragraphs.
    default_units = int(round(spc_bef_pt * 100))
    out: list[tuple[str, int | None]] = []
    slot_y_in = slots["y_positions_in"]
    prev_section_advance_pt = 0.0

    for i, section in enumerate(content[: slots["count"]]):
        section_paras = list(section) if section else [""]
        if i == 0:
            # No override; default spcBef on each paragraph
            for para_text in section_paras:
                out.append((para_text, default_units))
        else:
            target_gap_pt = (slot_y_in[i] - slot_y_in[i-1]) * 72.0
            override_pt = max(0.0, target_gap_pt - prev_section_advance_pt)
            override_units = int(round(override_pt * 100))
            out.append((section_paras[0], override_units))
            for para_text in section_paras[1:]:
                out.append((para_text, default_units))

        # Advance from this section's first-paragraph TOP to where the
        # next section's first paragraph would land if its override = 0.
        # The first paragraph's spcBef is excluded (absorbed by anchor).
        advance = 0.0
        for j, para_text in enumerate(section_paras):
            if j == 0:
                advance += visual_lines(para_text) * line_h_pt + spc_aft_pt
            else:
                advance += spc_bef_pt + visual_lines(para_text) * line_h_pt + spc_aft_pt
        prev_section_advance_pt = advance

    # First paragraph of section 0 keeps its inherited spcBef (set to None
    # so set_shape_text doesn't touch the template's original first-para
    # pPr, which already has the right spcBef).
    if out:
        out[0] = (out[0][0], None)
    return out


def _slot_aware_lines(shape, content: list[list[str]]) -> list[str]:
    """Legacy flat-strings shim: when caller doesn't apply spcBef
    overrides we fall back to concatenated paragraphs (no padding).
    Most callers use _slot_aware_paragraphs directly."""
    pairs = _slot_aware_paragraphs(shape, content)
    if not pairs:
        flat: list[str] = []
        for section in content:
            flat.extend(section)
        return flat or [""]
    return [text for text, _ in pairs] or [""]


def set_shape_text(shape, content: str | list[str] | list[list[str]]) -> None:
    """Replace the shape's text frame content while preserving font/style.

    - str → single paragraph
    - list[str] → one paragraph per item (bullets if the template paragraph
      already had bullet formatting at the paragraph-properties level)
    - list[list[str]] → slot mode: one inner list per slot anchor on the
      slide (numbered ovals/badges). Sections are padded with blank
      paragraphs so each visually aligns with its slot.
    """
    if not shape.has_text_frame:
        raise ValueError(f"Shape {shape.name!r} has no text frame")

    tf = shape.text_frame
    style_run = _style_run_from(tf)

    # pairs: list of (text, spc_bef_units | None). None inherits default;
    # an int sets an explicit `<a:spcBef><a:spcPts val=N/></a:spcBef>` on
    # that paragraph (units = pt × 100, OOXML format).
    pairs: list[tuple[str, int | None]]
    if isinstance(content, str):
        pairs = [(content, None)]
    elif content and isinstance(content[0], list):
        slot_pairs = _slot_aware_paragraphs(shape, content)  # type: ignore[arg-type]
        if slot_pairs:
            pairs = slot_pairs
        else:
            flat: list[str] = []
            for section in content:  # type: ignore[union-attr]
                flat.extend(section)
            pairs = [(t, None) for t in (flat or [""])]
    else:
        pairs = [(t, None) for t in (list(content) or [""])]  # type: ignore[arg-type]
    if not pairs:
        pairs = [("", None)]

    paragraphs = list(tf.paragraphs)
    first_text, first_spc = pairs[0]
    # Keep the first paragraph (preserves paragraph-properties: indent, bullet)
    _set_paragraph_text(paragraphs[0], first_text, style_run)
    if first_spc is not None:
        _set_paragraph_spc_bef(paragraphs[0]._p, first_spc)

    # Remove any other existing paragraphs after the first
    for p in paragraphs[1:]:
        p._p.getparent().remove(p._p)

    # Append additional paragraphs cloned from the first paragraph's pPr
    first_p = paragraphs[0]._p
    for extra_text, extra_spc in pairs[1:]:
        new_p = deepcopy(first_p)
        # strip all runs from clone, then add the line
        for r in new_p.findall(qn("a:r")):
            new_p.remove(r)
        first_p.addnext(new_p)
        from pptx.text.text import _Paragraph as ParaCls
        para_obj = ParaCls(new_p, tf)
        new_run = para_obj.add_run()
        new_run.text = extra_text
        if style_run is not None:
            _clone_run_style(style_run, new_run)
        # Section-start: explicit spcBef; otherwise strip any cloned override
        # so the paragraph inherits the frame's default spacing.
        _set_paragraph_spc_bef(new_p, extra_spc)
        first_p = new_p


def find_shapes_by_name(slide: Slide, name: str) -> list:
    """Some templates duplicate shape names (the IC template has many
    'Text Placeholder 2' shapes on slide 17, for example). Return all matches
    so the caller can disambiguate or set them all."""
    return [s for s in slide.shapes if s.name == name]


def set_text_by_name(
    slide: Slide, name: str, content: str | list[str] | list[list[str]],
) -> int:
    """Set text on every shape on `slide` whose name matches `name`.

    Returns the number of shapes that were updated. Returns 0 when the name
    does not exist on the slide (caller decides whether to warn).
    """
    matches = find_shapes_by_name(slide, name)
    for shape in matches:
        set_shape_text(shape, content)
    return len(matches)

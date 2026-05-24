"""Table cell injection and RAG-fill helpers."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.slide import Slide
from pptx.util import Pt

from ..models import RAG_HEX
from .anomalies import analyse_table
from .placeholders import _style_run_from, _set_paragraph_text


def find_table_by_name(slide: Slide, name: str):
    for shape in slide.shapes:
        if shape.name == name and shape.has_table:
            return shape.table
    return None


def _ensure_explicit_font_size(cell, fallback_pt: float | None) -> None:
    """Walk every run in the cell and set an explicit pt if missing.

    Without this, runs created in cells whose template counterpart had no
    explicit rPr/size inherit the table master's default — which renders
    inconsistently across PowerPoint, LibreOffice and Google Slides."""
    if fallback_pt is None:
        return
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is None:
                run.font.size = Pt(fallback_pt)


def set_table(slide: Slide, name: str, matrix: list[list[str]]) -> bool:
    """Apply a 2D matrix to the named table, top-left aligned.

    The template's existing dimensions are authoritative; the call also
    enforces a defensive explicit font size per cell using the table's body
    modal pt as fallback, so cells whose original was empty/divider don't
    inherit oversized master defaults at render time."""
    table = find_table_by_name(slide, name)
    if table is None:
        return False

    # Compute the table's body modal pt once — used to fix any cell whose
    # original had no explicit run-level font size.
    analysis = analyse_table(table)
    body_pt = analysis["body_modal_pt"]

    for r, row_values in enumerate(matrix):
        if r >= len(table.rows):
            break
        # Prefer this row's mode if present (e.g. header row has its own size);
        # else fall back to the table's body mode.
        row_schema = analysis["rows"][r] if r < len(analysis["rows"]) else None
        target_pt = (row_schema and row_schema["modal_pt"]) or body_pt
        for c, value in enumerate(row_values):
            if c >= len(table.columns):
                break
            cell = table.cell(r, c)
            tf = cell.text_frame
            style_run = _style_run_from(tf)
            paragraphs = list(tf.paragraphs)
            if paragraphs:
                _set_paragraph_text(paragraphs[0], value, style_run)
                for p in paragraphs[1:]:
                    p._p.getparent().remove(p._p)
            _ensure_explicit_font_size(cell, target_pt)
    return True


def set_cell_fill(slide: Slide, table_name: str, row: int, col: int, rag: str) -> bool:
    """Solid-fill a single table cell with a RAG color."""
    table = find_table_by_name(slide, table_name)
    if table is None:
        return False
    if row >= len(table.rows) or col >= len(table.columns):
        return False
    cell = table.cell(row, col)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor.from_string(RAG_HEX[rag])
    return True


def set_shape_rag_fill(slide: Slide, shape_name: str, rag: str) -> int:
    """Solid-fill every shape on the slide with the given name. Returns count."""
    count = 0
    for shape in slide.shapes:
        if shape.name != shape_name:
            continue
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(RAG_HEX[rag])
            count += 1
        except Exception:
            # Not every shape supports fill (graphic frames, connectors)
            continue
    return count


def apply_rag(slide: Slide, target: str, rag: str) -> int:
    """Dispatch a RAG update.

    `target` is either a shape name (rectangle / oval to fill), or a
    table cell reference of the form "TableName!r,c" (0-based)."""
    if "!" in target:
        table_name, coords = target.split("!", 1)
        r_str, c_str = coords.split(",")
        ok = set_cell_fill(slide, table_name, int(r_str), int(c_str), rag)
        return 1 if ok else 0
    return set_shape_rag_fill(slide, target, rag)
